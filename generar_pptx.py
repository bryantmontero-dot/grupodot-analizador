"""
Generador de PPTX comercial - Grupodot (Fase 3)
Toma el análisis estratégico de analisis_claude.py (Fase 2) y genera una
presentación PowerPoint con branding Grupodot, lista para enviar a un cliente.

Uso:
    python generar_pptx.py <ruta_analisis.json>
    python generar_pptx.py diagnosticos/analisis_www_rappi_com_co.json

Requiere:
    pip install python-pptx

El logo (logo-dot.png en la raíz) se inserta si existe; si no, se usa un
wordmark de texto "grupodot" con los colores de marca.

Autor: Equipo MarTech Grupodot
"""

import json
import os
import sys
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ============================================================
# CONFIGURACIÓN Y BRANDING GRUPODOT
# ============================================================

LOGO_PATH = "logo-dot.png"
DIR_DIAGNOSTICOS = "diagnosticos"

# Paleta Grupodot
AZUL = RGBColor(0x2E, 0xA6, 0xFF)       # Azul Digital
TURQUESA = RGBColor(0x38, 0xC0, 0xBD)   # Turquesa Digital
VERDE = RGBColor(0x48, 0xC5, 0x55)      # Verde Crecimiento
CORAL = RGBColor(0xF3, 0x5E, 0x59)      # Coral Decisión
AMARILLO = RGBColor(0xF3, 0xCB, 0x04)   # Amarillo Claridad
NEGRO = RGBColor(0x00, 0x00, 0x00)
BLANCO = RGBColor(0xFE, 0xFE, 0xFE)
GRIS = RGBColor(0x55, 0x55, 0x55)       # Texto secundario

# Mapeos semánticos
COLOR_IMPACTO = {"Alto": CORAL, "Medio": AMARILLO, "Bajo": VERDE}
ORDEN_IMPACTO = {"Alto": 0, "Medio": 1, "Bajo": 2}
COLOR_PRIORIDAD = {"Inmediata": CORAL, "Corto plazo": AMARILLO, "Mediano plazo": TURQUESA}
ORDEN_PRIORIDAD = {"Inmediata": 0, "Corto plazo": 1, "Mediano plazo": 2}

# Dimensiones (16:9)
ANCHO = Inches(13.333)
ALTO = Inches(7.5)


# ============================================================
# HELPERS DE CONSTRUCCIÓN
# ============================================================

def _slide_en_blanco(prs):
    """Añade y devuelve una slide con el layout en blanco."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fondo(slide, color):
    """Pinta el fondo completo de la slide con un color sólido."""
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, ALTO)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = color
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    # Enviar al fondo
    sp = fondo._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return fondo


def _caja_texto(slide, left, top, width, height, texto, size=18,
                color=NEGRO, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Calibri"):
    """Crea un textbox con formato y devuelve el text frame."""
    caja = slide.shapes.add_textbox(left, top, width, height)
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def _barra_titulo(slide, texto, color):
    """Dibuja la barra de cabecera de color con el título en blanco."""
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, Inches(1.1))
    barra.fill.solid()
    barra.fill.fore_color.rgb = color
    barra.line.fill.background()
    barra.shadow.inherit = False
    tf = barra.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLANCO
    run.font.name = "Calibri"
    return barra


def _pildora(slide, left, top, texto, color, width=Inches(1.5)):
    """Dibuja una etiqueta redondeada de color (impacto / prioridad)."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BLANCO
    run.font.name = "Calibri"
    return pill


def _insertar_logo_o_wordmark(slide, left, top, oscuro=True):
    """Inserta logo-dot.png si existe; si no, dibuja el wordmark 'grupodot'.

    Punto único donde se resuelve el logo faltante (se reemplazará el PNG después).
    """
    if os.path.isfile(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, left, top, height=Inches(0.8))
            return
        except Exception:
            pass  # si el PNG está corrupto, caemos al wordmark
    color_texto = NEGRO if oscuro else BLANCO
    tf = _caja_texto(slide, left, top, Inches(3), Inches(0.7), "grupo",
                     size=30, color=color_texto, bold=True)
    # Añadir "dot" en azul como segundo run en la misma línea
    run = tf.paragraphs[0].add_run()
    run.text = "dot"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = AZUL
    run.font.name = "Calibri"


# ============================================================
# SLIDES
# ============================================================

def _slide_portada(prs, analisis):
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    meta = analisis.get("meta", {})

    # Banda lateral de marca
    banda = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), ALTO)
    banda.fill.solid()
    banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    banda.shadow.inherit = False

    _insertar_logo_o_wordmark(slide, Inches(0.9), Inches(0.7), oscuro=True)

    _caja_texto(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2),
                "Diagnóstico MarTech", size=48, color=NEGRO, bold=True)
    _caja_texto(slide, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.6),
                "Análisis estratégico y oportunidades de crecimiento digital",
                size=20, color=GRIS)

    url = meta.get("url_analizada", "N/A")
    fecha = (meta.get("fecha_diagnostico") or "")[:10]
    _caja_texto(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.5),
                url, size=18, color=AZUL, bold=True)
    _caja_texto(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
                f"Grupodot · Análisis estratégico · {fecha}".strip(" ·"),
                size=12, color=GRIS)
    return slide


def _slide_resumen_madurez(prs, analisis):
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    _barra_titulo(slide, "Resumen ejecutivo", AZUL)

    madurez = analisis.get("madurez_digital", {})
    sector = analisis.get("sector", {})

    # Resumen ejecutivo (columna izquierda)
    _caja_texto(slide, Inches(0.6), Inches(1.5), Inches(7.6), Inches(4.5),
                analisis.get("resumen_ejecutivo", "N/A"), size=16, color=NEGRO)

    # Panel de madurez (columna derecha)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(8.5), Inches(1.5), Inches(4.2), Inches(4.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xFC)
    panel.line.color.rgb = AZUL
    panel.shadow.inherit = False

    nivel = madurez.get("nivel", "N/A")
    score = madurez.get("score", "N/A")
    _caja_texto(slide, Inches(8.7), Inches(1.7), Inches(3.8), Inches(0.5),
                "MADUREZ DIGITAL", size=13, color=GRIS, bold=True)
    _caja_texto(slide, Inches(8.7), Inches(2.2), Inches(3.8), Inches(0.8),
                str(nivel), size=30, color=AZUL, bold=True)
    _caja_texto(slide, Inches(8.7), Inches(3.1), Inches(3.8), Inches(0.6),
                f"Score: {score}/100", size=18, color=NEGRO, bold=True)

    # Barra de score
    try:
        pct = max(0, min(100, int(score))) / 100.0
    except (TypeError, ValueError):
        pct = 0
    riel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(8.7), Inches(3.7), Inches(3.8), Inches(0.25))
    riel.fill.solid()
    riel.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    riel.line.fill.background()
    riel.shadow.inherit = False
    if pct > 0:
        relleno = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(8.7), Inches(3.7), Inches(3.8 * pct), Inches(0.25))
        relleno.fill.solid()
        relleno.fill.fore_color.rgb = VERDE
        relleno.line.fill.background()
        relleno.shadow.inherit = False

    _caja_texto(slide, Inches(8.7), Inches(4.2), Inches(3.8), Inches(0.5),
                f"Sector: {sector.get('industria', 'N/A')}", size=13, color=NEGRO, bold=True)
    _caja_texto(slide, Inches(8.7), Inches(4.7), Inches(3.8), Inches(1.2),
                sector.get("subsector", ""), size=12, color=GRIS)
    return slide


def _slide_lista_simple(prs, titulo, color, items, fmt):
    """Slide genérica con una lista vertical de items formateados con `fmt(item)`."""
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    _barra_titulo(slide, titulo, color)

    top = Inches(1.5)
    for item in items:
        encabezado, detalle = fmt(item)
        # Viñeta de color
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top + Inches(0.08),
                                        Inches(0.18), Inches(0.18))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()
        bullet.shadow.inherit = False

        tf = _caja_texto(slide, Inches(1.0), top, Inches(11.7), Inches(0.9),
                         encabezado, size=16, color=NEGRO, bold=True)
        if detalle:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = detalle
            run.font.size = Pt(13)
            run.font.color.rgb = GRIS
            run.font.name = "Calibri"
        top += Inches(1.05)
    return slide


def _recortar(texto, max_chars):
    """Recorta un texto a `max_chars` añadiendo elipsis, sin cortar a media palabra."""
    texto = (texto or "").strip()
    if len(texto) <= max_chars:
        return texto
    recorte = texto[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:")
    return f"{recorte}…"


def _slide_oportunidades(prs, titulo, color_barra, oportunidades):
    """Slide de oportunidades filtradas: píldora de impacto + área en negrita +
    hallazgo corto (máx. 2 líneas) + recomendación en una línea con '→'."""
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    _barra_titulo(slide, titulo, color_barra)

    top = Inches(1.45)
    for o in oportunidades:
        impacto = o.get("impacto", "")
        color = COLOR_IMPACTO.get(impacto, GRIS)
        _pildora(slide, Inches(0.6), top + Inches(0.05), impacto or "—", color, width=Inches(1.2))

        # Área en negrita + hallazgo corto (no negrita) en el mismo párrafo.
        tf = _caja_texto(slide, Inches(2.0), top, Inches(10.7), Inches(0.9),
                         o.get("area", ""), size=15, color=NEGRO, bold=True)
        hallazgo = _recortar(o.get("hallazgo", ""), 150)
        if hallazgo:
            r = tf.paragraphs[0].add_run()
            r.text = f": {hallazgo}"
            r.font.size = Pt(14)
            r.font.bold = False
            r.font.color.rgb = NEGRO
            r.font.name = "Calibri"

        reco = o.get("recomendacion", "")
        if reco:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"→ {_recortar(reco, 130)}"
            run.font.size = Pt(12)
            run.font.color.rgb = GRIS
            run.font.name = "Calibri"
        top += Inches(1.0)
    return slide


def _slide_servicios(prs, analisis):
    servicios = sorted(
        analisis.get("servicios_recomendados", []),
        key=lambda s: ORDEN_PRIORIDAD.get(s.get("prioridad"), 9),
    )
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    _barra_titulo(slide, "Servicios recomendados", TURQUESA)

    top = Inches(1.45)
    for s in servicios:
        prioridad = s.get("prioridad", "")
        color = COLOR_PRIORIDAD.get(prioridad, GRIS)
        _pildora(slide, Inches(0.6), top + Inches(0.05), prioridad or "—", color, width=Inches(1.7))
        tf = _caja_texto(slide, Inches(2.5), top, Inches(10.2), Inches(0.9),
                         s.get("servicio", ""), size=15, color=NEGRO, bold=True)
        desc = s.get("descripcion", "")
        if desc:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = desc
            run.font.size = Pt(12)
            run.font.color.rgb = GRIS
            run.font.name = "Calibri"
        top += Inches(1.05)
    return slide


def _slide_quick_wins(prs, analisis):
    slide = _slide_en_blanco(prs)
    _fondo(slide, BLANCO)
    _barra_titulo(slide, "Quick wins", VERDE)

    top = Inches(1.6)
    for i, qw in enumerate(analisis.get("quick_wins", []), 1):
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top,
                                     Inches(0.5), Inches(0.5))
        num.fill.solid()
        num.fill.fore_color.rgb = VERDE
        num.line.fill.background()
        num.shadow.inherit = False
        ptf = num.text_frame
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run()
        r.text = str(i)
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = BLANCO
        r.font.name = "Calibri"

        _caja_texto(slide, Inches(1.3), top, Inches(11.4), Inches(0.7),
                    qw, size=16, color=NEGRO, anchor=MSO_ANCHOR.MIDDLE)
        top += Inches(0.85)
    return slide


def _slide_cierre(prs, analisis):
    slide = _slide_en_blanco(prs)
    _fondo(slide, NEGRO)
    _insertar_logo_o_wordmark(slide, Inches(0.9), Inches(0.7), oscuro=False)

    _caja_texto(slide, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.5),
                "Hablemos de tu estrategia MarTech", size=40, color=BLANCO, bold=True)
    _caja_texto(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
                "Convirtamos estas oportunidades en crecimiento medible.",
                size=18, color=TURQUESA)
    _caja_texto(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
                "grupodot.com  ·  contacto@grupodot.com", size=14, color=BLANCO)
    return slide


# ============================================================
# ORQUESTADOR
# ============================================================

def construir_presentacion(analisis: dict) -> Presentation:
    """Construye la presentación completa a partir del análisis de Fase 2."""
    prs = Presentation()
    prs.slide_width = ANCHO
    prs.slide_height = ALTO

    # Slides 1-3
    _slide_portada(prs, analisis)
    _slide_resumen_madurez(prs, analisis)
    _slide_lista_simple(
        prs, "Fortalezas", VERDE, analisis.get("fortalezas", []),
        lambda f: (f.get("area", ""), f.get("detalle", "")),
    )

    # Slides 4-5: oportunidades divididas por impacto (se omite la slide vacía).
    oportunidades = analisis.get("oportunidades", [])
    criticas = [o for o in oportunidades if o.get("impacto") == "Alto"]
    mejora = sorted(
        [o for o in oportunidades if o.get("impacto") in ("Medio", "Bajo")],
        key=lambda o: ORDEN_IMPACTO.get(o.get("impacto"), 9),
    )
    if criticas:
        _slide_oportunidades(prs, "Oportunidades críticas", CORAL, criticas)
    if mejora:
        _slide_oportunidades(prs, "Oportunidades de mejora", AMARILLO, mejora)

    # Slides 6-8
    _slide_servicios(prs, analisis)
    _slide_quick_wins(prs, analisis)
    _slide_cierre(prs, analisis)

    return prs


# ============================================================
# ENTRY POINT
# ============================================================

def main():
    if len(sys.argv) < 2:
        print("Uso: python generar_pptx.py <ruta_analisis.json>")
        print("Ejemplo: python generar_pptx.py diagnosticos/analisis_www_rappi_com_co.json")
        raise SystemExit(1)

    archivo_analisis = sys.argv[1]

    if not os.path.isfile(archivo_analisis):
        print(f"ERROR: No se encontró el archivo de análisis: {archivo_analisis}")
        raise SystemExit(1)

    print(f"Cargando análisis: {archivo_analisis}")
    with open(archivo_analisis, "r", encoding="utf-8") as f:
        analisis = json.load(f)

    print("Construyendo presentación...")
    prs = construir_presentacion(analisis)

    # Guardar en diagnosticos/, dominio derivado de la URL analizada.
    os.makedirs(DIR_DIAGNOSTICOS, exist_ok=True)
    url = analisis.get("meta", {}).get("url_analizada", "sitio")
    dominio = urlparse(url).netloc.replace(".", "_") or "sitio"
    output_file = os.path.join(DIR_DIAGNOSTICOS, f"presentacion_{dominio}.pptx")
    prs.save(output_file)

    print(f"\nPresentación guardada en: {output_file}")


if __name__ == "__main__":
    main()
